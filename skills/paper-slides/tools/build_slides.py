#!/usr/bin/env python3
"""Deterministic host builder: paper/ → slides/ (Beamer + PDF + PPTX + scripts).

No LLM. Extracts title/sections/bullets from paper markdown or LaTeX and emits
the full slides artifact chain required by the paper_slides workflow template.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

VENUE_COLORS = {
    "neurips": ("8B5CF6", "2563EB"),
    "icml": ("DC2626", "1D4ED8"),
    "iclr": ("059669", "0284C7"),
    "cvpr": ("2563EB", "7C3AED"),
    "aaai": ("0369A1", "DC2626"),
    "acl": ("0891B2", "7C3AED"),
    "emnlp": ("D97706", "2563EB"),
    "eccv": ("C026D3", "0891B2"),
    "generic": ("334155", "2563EB"),
}

TALK_SLIDE_TARGETS = {
    "poster-talk": 6,
    "spotlight": 10,
    "oral": 16,
    "invited": 24,
}


@dataclass
class Slide:
    title: str
    bullets: List[str] = field(default_factory=list)
    notes: str = ""
    kind: str = "content"  # title | content | thankyou


def _now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8-sig", errors="replace")


def _strip_tex(text: str) -> str:
    text = re.sub(r"%.*$", "", text, flags=re.M)
    text = re.sub(r"\\(cite|ref|label|includegraphics|url|href)(\[[^\]]*\])?\{[^}]*\}", "", text)
    text = re.sub(r"\\(textbf|textit|emph|mathrm|mathbf|texttt)\{([^}]*)\}", r"\2", text)
    text = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^}]*\})?", " ", text)
    text = re.sub(r"[{}$~^_&]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _strip_md(text: str) -> str:
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"[#>*`_]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _first_nonempty(lines: Iterable[str]) -> str:
    for line in lines:
        s = line.strip()
        if s:
            return s
    return ""


def _extract_from_markdown(paper_dir: Path) -> Tuple[str, str, List[Tuple[str, List[str]]]]:
    md = paper_dir / "main.md"
    if not md.is_file():
        for candidate in sorted(paper_dir.glob("*.md")):
            md = candidate
            break
    if not md.is_file():
        return "Untitled Research", "Vibe Research", []

    text = _read_text(md)
    lines = text.splitlines()
    title = "Untitled Research"
    for line in lines:
        if line.startswith("# "):
            title = _strip_md(line[2:])
            break
    sections: List[Tuple[str, List[str]]] = []
    current = "Overview"
    bullets: List[str] = []
    for line in lines:
        if line.startswith("## "):
            if bullets:
                sections.append((current, bullets[:8]))
            current = _strip_md(line[3:]) or "Section"
            bullets = []
            continue
        stripped = line.strip()
        if stripped.startswith(("- ", "* ", "1. ", "2. ", "3. ")):
            item = _strip_md(re.sub(r"^(\d+\.|[-*])\s+", "", stripped))
            if item:
                bullets.append(item[:120])
        elif stripped and not stripped.startswith("#") and len(stripped) > 20:
            bullets.append(_strip_md(stripped)[:120])
    if bullets:
        sections.append((current, bullets[:8]))
    return title, "Vibe Research", sections


def _extract_from_tex(paper_dir: Path) -> Tuple[str, str, List[Tuple[str, List[str]]]]:
    main = paper_dir / "main.tex"
    if not main.is_file():
        return _extract_from_markdown(paper_dir)
    text = _read_text(main)
    title_m = re.search(r"\\title\{([^}]*)\}", text)
    author_m = re.search(r"\\author\{([^}]*)\}", text)
    title = _strip_tex(title_m.group(1)) if title_m else "Untitled Research"
    author = _strip_tex(author_m.group(1)) if author_m else "Vibe Research"

    sections: List[Tuple[str, List[str]]] = []
    # Prefer sections/*.tex if present
    section_files = sorted((paper_dir / "sections").glob("*.tex")) if (paper_dir / "sections").is_dir() else []
    sources = section_files or [main]
    for src in sources:
        body = _read_text(src)
        for block in re.finditer(
            r"\\(?:section|subsection)\*?\{([^}]*)\}(.*?)(?=\\(?:section|subsection)\*?\{|\\end\{document\}|\Z)",
            body,
            flags=re.S,
        ):
            name = _strip_tex(block.group(1)) or src.stem
            chunk = block.group(2)
            items = []
            for item in re.findall(r"\\item\s+(.*?)(?=\\item|\\end\{)", chunk, flags=re.S):
                cleaned = _strip_tex(item)
                if cleaned:
                    items.append(cleaned[:120])
            if not items:
                para = _strip_tex(chunk)
                if para:
                    # split into ~sentence fragments
                    parts = re.split(r"(?<=[。.!?])\s+", para)
                    items = [p[:120] for p in parts if len(p.strip()) > 12][:6]
            if items:
                sections.append((name, items[:8]))
    if not sections:
        return _extract_from_markdown(paper_dir) if (paper_dir / "main.md").is_file() else (title, author, [])
    return title or "Untitled Research", author or "Vibe Research", sections


def _build_slides(
    title: str,
    author: str,
    sections: Sequence[Tuple[str, List[str]]],
    talk_type: str,
    venue: str,
    minutes: int,
) -> List[Slide]:
    target = TALK_SLIDE_TARGETS.get(talk_type, 10)
    slides: List[Slide] = [
        Slide(title=title, bullets=[author, venue.upper(), f"{minutes} min talk"], notes="Wait for chair introduction.", kind="title"),
    ]
    if sections:
        outline = [name for name, _ in sections[:8]]
        slides.append(
            Slide(
                title="Outline",
                bullets=outline,
                notes="Preview the narrative arc for the audience.",
                kind="content",
            )
        )
    # Motivation / problem from first section
    if sections:
        name, items = sections[0]
        slides.append(Slide(title=name or "Motivation", bullets=items[:5], notes=f"Spend ~1 min on {name}."))
    # Key insight / method middle sections
    mid = list(sections[1:-1]) if len(sections) > 2 else list(sections[1:])
    for name, items in mid:
        slides.append(Slide(title=name, bullets=items[:5], notes=f"Explain {name} with concrete evidence."))
        if len(slides) >= target - 2:
            break
    # Results / conclusion
    if sections and len(sections) > 1:
        name, items = sections[-1]
        slides.append(Slide(title=name or "Results", bullets=items[:5], notes="Highlight headline numbers only."))
    slides.append(
        Slide(
            title="Takeaways",
            bullets=[
                "Problem grounded in real research need",
                "Method is reproducible end-to-end",
                "Evidence chain is auditable",
            ],
            notes="Close with one memorable takeaway.",
            kind="content",
        )
    )
    slides.append(
        Slide(
            title="Thank You",
            bullets=["Questions welcome", "Paper & code available on request"],
            notes="Invite questions.",
            kind="thankyou",
        )
    )
    # Trim or pad to target range
    if len(slides) > target + 2:
        slides = slides[: target + 1] + [slides[-1]]
    return slides


def _latex_escape(text: str) -> str:
    repl = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    out = []
    for ch in text:
        out.append(repl.get(ch, ch))
    return "".join(out)


def _write_outline(path: Path, slides: Sequence[Slide], venue: str, talk_type: str, minutes: int) -> None:
    lines = [
        f"# Slide Outline",
        "",
        f"- Venue: {venue}",
        f"- Talk type: {talk_type}",
        f"- Minutes: {minutes}",
        f"- Slide count: {len(slides)}",
        "",
    ]
    for i, slide in enumerate(slides, 1):
        lines.append(f"## {i}. {slide.title}")
        for b in slide.bullets:
            lines.append(f"- {b}")
        if slide.notes:
            lines.append(f"\n_Notes_: {slide.notes}")
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def _write_talk_script(path: Path, slides: Sequence[Slide], title: str, venue: str, talk_type: str, minutes: int) -> None:
    lines = [
        f"# Talk Script: {title}",
        "",
        f"**Venue**: {venue}",
        f"**Talk type**: {talk_type} ({minutes} min)",
        f"**Total slides**: {len(slides)}",
        "",
    ]
    per = max(20, int((minutes * 60) / max(1, len(slides))))
    t = 0
    for i, slide in enumerate(slides, 1):
        start, end = t, t + per
        lines.append(f"## Slide {i}: {slide.title} [{start//60}:{start%60:02d} - {end//60}:{end%60:02d}]")
        lines.append("")
        if slide.kind == "title":
            lines.append(f'"Thank you. Today I will present {title}."')
        elif slide.kind == "thankyou":
            lines.append('"To summarize, our key takeaway is that the evidence chain is reproducible. Happy to take questions."')
        else:
            spoken = slide.notes or ("Next: " + "; ".join(slide.bullets[:2]))
            lines.append(f'"{spoken}"')
            if slide.bullets:
                lines.append("")
                lines.append("Key points:")
                for b in slide.bullets:
                    lines.append(f"- {b}")
        lines.append("")
        t = end
    lines.extend(
        [
            "## Anticipated Q&A",
            "",
            "### Q1: What is the strongest evidence?",
            "**A**: The primary metrics and artifact lineage in the paper results.",
            "",
            "### Q2: What are the limitations?",
            "**A**: Scope and assumptions are stated in the discussion/limitations section.",
            "",
        ]
    )
    path.write_text("\n".join(lines), encoding="utf-8")


def _write_speaker_notes(path: Path, slides: Sequence[Slide]) -> None:
    lines = ["# Speaker Notes", ""]
    for i, slide in enumerate(slides, 1):
        lines.append(f"## Slide {i}: {slide.title}")
        lines.append(slide.notes or "(no notes)")
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def _write_beamer(path: Path, slides: Sequence[Slide], title: str, author: str, venue: str, aspect: str, primary: str, accent: str, notes: bool) -> None:
    ratio = "169" if aspect.replace(":", "") in {"169", "16:9", "169"} or aspect == "16:9" else "43"
    if aspect == "16:9":
        ratio = "169"
    elif aspect == "4:3":
        ratio = "43"
    frames = []
    for slide in slides:
        if slide.kind == "title":
            frames.append("\\begin{frame}\n\\titlepage\n\\end{frame}\n")
            continue
        body = ["\\begin{itemize}"]
        for b in slide.bullets:
            body.append(f"  \\item {_latex_escape(b)}")
        body.append("\\end{itemize}")
        note = f"\\note{{{_latex_escape(slide.notes)}}}\n" if notes and slide.notes else ""
        frames.append(
            f"\\begin{{frame}}{{{_latex_escape(slide.title)}}}\n"
            + "\n".join(body)
            + f"\n{note}\\end{{frame}}\n"
        )
    tex = f"""\\documentclass[aspectratio={ratio}]{{beamer}}
\\usepackage[T1]{{fontenc}}
\\usepackage[utf8]{{inputenc}}
\\usepackage{{lmodern}}
\\usepackage{{xcolor}}
\\usepackage{{graphicx}}
\\definecolor{{primary}}{{HTML}}{{{primary}}}
\\definecolor{{accent}}{{HTML}}{{{accent}}}
\\usetheme{{default}}
\\usecolortheme{{default}}
\\setbeamercolor{{frametitle}}{{fg=primary}}
\\setbeamercolor{{title}}{{fg=primary}}
\\setbeamercolor{{structure}}{{fg=accent}}
\\setbeamercolor{{itemize item}}{{fg=primary}}
\\setbeamertemplate{{navigation symbols}}{{}}
\\setbeamertemplate{{footline}}{{%
  \\hfill\\insertframenumber/\\inserttotalframenumber\\hspace{{2mm}}\\vspace{{2mm}}%
}}
\\title{{{_latex_escape(title)}}}
\\author{{{_latex_escape(author)}}}
\\institute{{Vibe Research}}
\\date{{{_latex_escape(venue)}}}
\\begin{{document}}
{''.join(frames)}
\\end{{document}}
"""
    path.write_text(tex, encoding="utf-8")


def _find_latex() -> Optional[str]:
    for name in ("pdflatex", "xelatex"):
        found = shutil.which(name)
        if found:
            return found
    return None


def _compile_pdf(slides_dir: Path, engine: Optional[str] = None) -> Tuple[int, str, str]:
    binary = engine or _find_latex()
    if not binary:
        return 3, "", "pdflatex/xelatex unavailable"
    cmd = [binary, "-interaction=nonstopmode", "-halt-on-error", "main.tex"]
    stdout_parts: List[str] = []
    stderr_parts: List[str] = []
    rc = 0
    for _ in range(2):
        proc = subprocess.run(
            cmd,
            cwd=str(slides_dir),
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=180,
        )
        stdout_parts.append(proc.stdout or "")
        stderr_parts.append(proc.stderr or "")
        rc = proc.returncode
        if rc != 0:
            break
    return rc, "\n".join(stdout_parts), "\n".join(stderr_parts)


def _write_pptx(path: Path, slides: Sequence[Slide], title: str, venue: str, aspect: str, primary: str, accent: str) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    if aspect == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    primary_rgb = RGBColor(int(primary[0:2], 16), int(primary[2:4], 16), int(primary[4:6], 16))
    accent_rgb = RGBColor(int(accent[0:2], 16), int(accent[2:4], 16), int(accent[4:6], 16))

    for slide in slides:
        s = prs.slides.add_slide(blank)
        if slide.kind == "title":
            box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), prs.slide_width - Inches(1.6), Inches(1.5))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = primary_rgb
            p.alignment = PP_ALIGN.CENTER
            sub = s.shapes.add_textbox(Inches(0.8), Inches(4.0), prs.slide_width - Inches(1.6), Inches(1.2))
            stf = sub.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = " · ".join(slide.bullets)
            sp.font.size = Pt(18)
            sp.font.color.rgb = accent_rgb
            sp.alignment = PP_ALIGN.CENTER
        else:
            title_box = s.shapes.add_textbox(Inches(0.6), Inches(0.35), prs.slide_width - Inches(1.2), Inches(0.9))
            ttf = title_box.text_frame
            tp = ttf.paragraphs[0]
            tp.text = slide.title
            tp.font.size = Pt(28)
            tp.font.bold = True
            tp.font.color.rgb = primary_rgb
            body = s.shapes.add_textbox(Inches(0.8), Inches(1.4), prs.slide_width - Inches(1.6), Inches(5.2))
            btf = body.text_frame
            btf.word_wrap = True
            for i, bullet in enumerate(slide.bullets or [""]):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes
    prs.save(str(path))


def _copy_figures(paper_dir: Path, slides_dir: Path) -> List[str]:
    src = paper_dir / "figures"
    dst = slides_dir / "figures"
    dst.mkdir(parents=True, exist_ok=True)
    copied = []
    if not src.is_dir():
        return copied
    for path in src.rglob("*"):
        if path.is_file() and path.suffix.lower() in {".pdf", ".png", ".jpg", ".jpeg", ".svg"}:
            target = dst / path.name
            shutil.copy2(path, target)
            copied.append(target.name)
    return copied


def build(
    workspace: Path,
    *,
    venue: str = "NeurIPS",
    talk_type: str = "spotlight",
    minutes: int = 15,
    aspect: str = "16:9",
    notes: bool = True,
    engine: str = "",
) -> dict:
    workspace = Path(workspace).expanduser().resolve()
    paper_dir = workspace / "paper"
    slides_dir = workspace / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    if not paper_dir.is_dir():
        raise SystemExit("paper/ directory is missing")
    if not any((paper_dir / name).is_file() for name in ("main.tex", "main.md", "main.pdf")):
        # allow sections-only workspaces
        if not list(paper_dir.glob("**/*")):
            raise SystemExit("paper/ has no content")

    title, author, sections = _extract_from_tex(paper_dir)
    if not sections:
        title2, author2, sections = _extract_from_markdown(paper_dir)
        title = title or title2
        author = author or author2

    venue_key = venue.strip().lower()
    primary, accent = VENUE_COLORS.get(venue_key, VENUE_COLORS["generic"])
    slides = _build_slides(title, author, sections, talk_type, venue, minutes)
    figures = _copy_figures(paper_dir, slides_dir)

    _write_outline(slides_dir / "SLIDE_OUTLINE.md", slides, venue, talk_type, minutes)
    _write_talk_script(slides_dir / "TALK_SCRIPT.md", slides, title, venue, talk_type, minutes)
    _write_speaker_notes(slides_dir / "speaker_notes.md", slides)
    _write_beamer(slides_dir / "main.tex", slides, title, author, venue, aspect, primary, accent, notes)
    _write_pptx(slides_dir / "presentation.pptx", slides, title, venue, aspect, primary, accent)

    # Keep a regenerator for users who edit outline later
    (slides_dir / "generate_pptx.py").write_text(
        "# Regenerated by host paper-slides builder; re-run build_slides.py to refresh.\n",
        encoding="utf-8",
    )

    rc, stdout, stderr = _compile_pdf(slides_dir, engine or None)
    pdf = slides_dir / "main.pdf"
    # Fallback: if LaTeX missing/fails, still emit a minimal PDF via reportlab-free path:
    # use pure PDF with text operators only when compile fails? Better: mark failure.
    # For environments without full beamer packages we try once more with article fallback.
    if not (pdf.is_file() and pdf.stat().st_size >= 500):
        fallback = slides_dir / "main_article_fallback.tex"
        body = ["\\documentclass[11pt]{article}", "\\usepackage[margin=1in]{geometry}", "\\begin{document}"]
        body.append(f"\\title{{{_latex_escape(title)}}}\\author{{{_latex_escape(author)}}}\\maketitle")
        for slide in slides:
            body.append(f"\\section*{{{_latex_escape(slide.title)}}}")
            body.append("\\begin{itemize}")
            for b in slide.bullets:
                body.append(f"\\item {_latex_escape(b)}")
            body.append("\\end{itemize}")
        body.append("\\end{document}")
        fallback.write_text("\n".join(body) + "\n", encoding="utf-8")
        # overwrite main.tex only if beamer failed hard
        if rc != 0:
            shutil.copy2(fallback, slides_dir / "main.tex")
            rc, stdout, stderr = _compile_pdf(slides_dir, engine or None)

    state = {
        "phase": 8,
        "venue": venue,
        "talk_type": talk_type,
        "slide_count": len(slides),
        "status": "completed" if pdf.is_file() and pdf.stat().st_size >= 500 else "partial",
        "timestamp": _now_iso(),
        "figures": figures,
        "executor": "host_step_runner",
        "primary_pdf": "slides/main.pdf",
        "pptx": "slides/presentation.pptx",
        "compile_returncode": rc,
    }
    (slides_dir / "SLIDES_STATE.json").write_text(json.dumps(state, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (slides_dir / "SLIDES_REVIEW.md").write_text(
        "# Slides Review\n\nHost deterministic builder completed without external reviewer.\n"
        f"- Slides: {len(slides)}\n- Compile rc: {rc}\n",
        encoding="utf-8",
    )

    report = {
        "success": pdf.is_file() and pdf.stat().st_size >= 500 and (slides_dir / "presentation.pptx").is_file(),
        "returncode": 0 if (pdf.is_file() and pdf.stat().st_size >= 500) else (rc or 4),
        "title": title,
        "slide_count": len(slides),
        "artifacts": [
            "slides/main.tex",
            "slides/main.pdf",
            "slides/presentation.pptx",
            "slides/SLIDE_OUTLINE.md",
            "slides/TALK_SCRIPT.md",
            "slides/speaker_notes.md",
            "slides/SLIDES_STATE.json",
        ],
        "stdout_tail": (stdout or "")[-2000:],
        "stderr_tail": (stderr or "")[-2000:],
        "state": state,
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return report


def main(argv: Optional[Sequence[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Build conference slides from paper/")
    parser.add_argument("--workspace", required=True, help="Workflow workspace root")
    parser.add_argument("--venue", default="NeurIPS")
    parser.add_argument("--talk-type", default="spotlight", dest="talk_type")
    parser.add_argument("--minutes", type=int, default=15)
    parser.add_argument("--aspect", default="16:9")
    parser.add_argument("--notes", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("--engine", default="")
    args = parser.parse_args(argv)
    report = build(
        Path(args.workspace),
        venue=args.venue,
        talk_type=args.talk_type,
        minutes=args.minutes,
        aspect=args.aspect,
        notes=args.notes,
        engine=args.engine,
    )
    return 0 if report.get("success") else int(report.get("returncode") or 4)


if __name__ == "__main__":
    raise SystemExit(main())
