#!/usr/bin/env python3
"""Deterministic host builder: paper/ → poster/ (LaTeX + PDF + PPTX + speech).

No LLM. Extracts title/sections from paper markdown or LaTeX and emits the
paper_poster workflow artifact chain. Uses article+geometry multi-column layout
(portable; does not require tcbposter).
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional, Sequence, Tuple

VENUE_THEME = {
    "neurips": ("4C1D95", "6D28D9", "2563EB", "F5F3FF"),
    "icml": ("7F1D1D", "B91C1C", "1E40AF", "FDF2F2"),
    "iclr": ("065F46", "059669", "0284C7", "F0FDF4"),
    "cvpr": ("1E3A8A", "2563EB", "7C3AED", "F8FAFC"),
    "aaai": ("0C4A6E", "0369A1", "DC2626", "F0F9FF"),
    "acl": ("155E75", "0891B2", "7C3AED", "F0FDFA"),
    "emnlp": ("713F12", "D97706", "2563EB", "FFFBEB"),
    "eccv": ("701A75", "C026D3", "0891B2", "FDF4FF"),
    "generic": ("1E293B", "334155", "2563EB", "F8FAFC"),
}

SIZE_MM = {
    ("A0", "landscape"): (1189, 841),
    ("A0", "portrait"): (841, 1189),
    ("A1", "landscape"): (841, 594),
    ("A1", "portrait"): (594, 841),
}


@dataclass
class PosterBox:
    title: str
    bullets: List[str] = field(default_factory=list)


def _now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8-sig", errors="replace")


def _strip_tex(text: str) -> str:
    text = re.sub(r"%.*$", "", text, flags=re.M)
    text = re.sub(r"\\(cite|ref|label|includegraphics|url|href)(\[[^\]]*\])?\{[^}]*\}", "", text)
    text = re.sub(r"\\(textbf|textit|emph|mathrm|mathbf|texttt)\{([^}]*)\}", r"\2", text)
    text = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^}]*\})?", " ", text)
    text = re.sub(r"[{}$~^_&]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _strip_md(text: str) -> str:
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"[#>*`_]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _latex_escape(text: str) -> str:
    repl = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    return "".join(repl.get(ch, ch) for ch in text)


def _extract(paper_dir: Path) -> Tuple[str, str, List[Tuple[str, List[str]]], List[str]]:
    title = "Untitled Research"
    author = "Vibe Research"
    sections: List[Tuple[str, List[str]]] = []
    stats: List[str] = []

    main_tex = paper_dir / "main.tex"
    main_md = paper_dir / "main.md"
    if main_tex.is_file():
        text = _read_text(main_tex)
        tm = re.search(r"\\title\{([^}]*)\}", text)
        am = re.search(r"\\author\{([^}]*)\}", text)
        if tm:
            title = _strip_tex(tm.group(1)) or title
        if am:
            author = _strip_tex(am.group(1)) or author
        sources = sorted((paper_dir / "sections").glob("*.tex")) if (paper_dir / "sections").is_dir() else [main_tex]
        for src in sources:
            body = _read_text(src)
            for block in re.finditer(
                r"\\(?:section|subsection)\*?\{([^}]*)\}(.*?)(?=\\(?:section|subsection)\*?\{|\\end\{document\}|\Z)",
                body,
                flags=re.S,
            ):
                name = _strip_tex(block.group(1)) or src.stem
                chunk = _strip_tex(block.group(2))
                items = [p[:100] for p in re.split(r"(?<=[。.!?])\s+", chunk) if len(p.strip()) > 12][:5]
                if items:
                    sections.append((name, items))
                for m in re.findall(r"(\d+(?:\.\d+)?\s*%|\d+(?:\.\d+)?x|\bF1\b|\bAUC\b)", chunk, flags=re.I):
                    stats.append(m)
    if main_md.is_file() and not sections:
        text = _read_text(main_md)
        for line in text.splitlines():
            if line.startswith("# "):
                title = _strip_md(line[2:]) or title
                break
        current = "Overview"
        bullets: List[str] = []
        for line in text.splitlines():
            if line.startswith("## "):
                if bullets:
                    sections.append((current, bullets[:5]))
                current = _strip_md(line[3:]) or "Section"
                bullets = []
                continue
            s = line.strip()
            if s.startswith(("- ", "* ")):
                bullets.append(_strip_md(s[2:])[:100])
            elif s and not s.startswith("#") and len(s) > 20:
                bullets.append(_strip_md(s)[:100])
            for m in re.findall(r"(\d+(?:\.\d+)?\s*%|\d+(?:\.\d+)?x)", s):
                stats.append(m)
        if bullets:
            sections.append((current, bullets[:5]))
    # unique stats
    seen = set()
    uniq_stats = []
    for s in stats:
        if s not in seen:
            seen.add(s)
            uniq_stats.append(s)
    if not uniq_stats:
        uniq_stats = ["E2E", "Local", "Auditable", "Reproducible"]
    return title, author, sections, uniq_stats[:4]


def _plan_boxes(sections: Sequence[Tuple[str, List[str]]]) -> List[PosterBox]:
    defaults = [
        ("Motivation", ["Problem is real and costly", "Prior work leaves a gap", "Need a verifiable pipeline"]),
        ("Method", ["Deterministic host builders", "Evidence-native gates", "Multi-provider agents"]),
        ("Results", ["Artifacts compile locally", "Lineage recorded per step", "No silent degradation"]),
        ("Takeaways", ["Doctoral workflow coverage", "Brand-clean runtime", "UI→API→executor→artifact"]),
    ]
    boxes: List[PosterBox] = []
    for i, (name, items) in enumerate(sections[:6]):
        boxes.append(PosterBox(title=name, bullets=items[:4] or defaults[min(i, len(defaults) - 1)][1]))
    while len(boxes) < 4:
        name, items = defaults[len(boxes)]
        boxes.append(PosterBox(title=name, bullets=items))
    return boxes[:6]


def _write_plan(path: Path, title: str, venue: str, size: str, orientation: str, columns: int, boxes: Sequence[PosterBox], stats: Sequence[str]) -> None:
    lines = [
        "# Poster Content Plan",
        "",
        f"- Title: {title}",
        f"- Venue: {venue}",
        f"- Size: {size} {orientation}",
        f"- Columns: {columns}",
        f"- Stats: {', '.join(stats)}",
        "",
    ]
    for i, box in enumerate(boxes, 1):
        lines.append(f"## Box {i}: {box.title}")
        for b in box.bullets:
            lines.append(f"- {b}")
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def _write_speech(path: Path, title: str, venue: str, boxes: Sequence[PosterBox]) -> None:
    lines = [
        "# Poster Presentation Script",
        "",
        f"**Paper**: {title}",
        f"**Venue**: {venue}",
        "**Estimated time**: 2-3 minutes",
        "",
        "## Opening (15 seconds)",
        f'"Thanks for stopping by! This poster presents {title}."',
        "",
    ]
    for box in boxes[:4]:
        lines.append(f"## {box.title}")
        if box.bullets:
            lines.append(f'"{box.bullets[0]}"')
        lines.append("")
    lines.extend(
        [
            "## Closing",
            '"Happy to discuss details — paper and code are available on request."',
            "",
            "## Anticipated Q&A",
            "",
            "### Q1: How do you ensure reproducibility?",
            "**A**: Host builders write lineage JSON and real artifacts under the workspace.",
            "",
        ]
    )
    path.write_text("\n".join(lines), encoding="utf-8")


def _write_tex(
    path: Path,
    title: str,
    author: str,
    venue: str,
    size: str,
    orientation: str,
    columns: int,
    boxes: Sequence[PosterBox],
    stats: Sequence[str],
    primary: str,
    secondary: str,
    accent: str,
    bg: str,
) -> None:
    w, h = SIZE_MM[(size.upper(), orientation.lower())]
    # body font scales down for A1
    body_pt = 28 if size.upper() == "A0" else 18
    title_pt = 48 if size.upper() == "A0" else 32
    stat_pt = 36 if size.upper() == "A0" else 22

    # Use fixed fractional widths (portable; avoids dimexpr * unit pitfalls).
    col_frac = {
        2: "0.48",
        3: "0.31",
        4: "0.23",
    }.get(columns, "0.31")
    box_chunks = []
    for box in boxes:
        items = "\n".join(f"    \\item {_latex_escape(b)}" for b in box.bullets)
        box_chunks.append(
            f"""\\begin{{minipage}}[t]{{{col_frac}\\linewidth}}
\\vspace{{0pt}}
\\colorbox{{white}}{{%
\\parbox[t]{{0.96\\linewidth}}{{%
  \\textcolor{{secondary}}{{\\textbf{{\\large {_latex_escape(box.title)}}}}}\\par\\vspace{{4mm}}
  \\begin{{itemize}}\\setlength\\itemsep{{3mm}}
{items}
  \\end{{itemize}}
}}}}
\\end{{minipage}}"""
        )
    # Wrap rows of `columns` minipages
    rows = []
    row: List[str] = []
    for chunk in box_chunks:
        row.append(chunk)
        if len(row) >= columns:
            rows.append("\\hfill\n".join(row))
            row = []
    if row:
        rows.append("\\hfill\n".join(row))
    body = "\n\\vspace{8mm}\n\\noindent\n".join(rows) if rows else ""
    stat_parts = []
    for s in (list(stats) + ["—"] * 4)[:4]:
        escaped = _latex_escape(s)
        stat_parts.append(
            "\\fcolorbox{primary}{white}{\\parbox{0.2\\linewidth}{\\centering\\vspace{2mm}"
            f"{{\\fontsize{{{stat_pt}}}{{{stat_pt + 8}}}\\selectfont\\bfseries"
            f"\\textcolor{{primary}}{{{escaped}}}}}"
            "\\vspace{2mm}}}"
        )
    stats_row = " \\hfill ".join(stat_parts)

    tex = f"""\\documentclass{{article}}
\\usepackage[paperwidth={w}mm,paperheight={h}mm,margin=12mm]{{geometry}}
\\usepackage[table]{{xcolor}}
\\usepackage{{graphicx}}
\\usepackage{{enumitem}}
\\usepackage{{lmodern}}
\\usepackage[T1]{{fontenc}}
\\usepackage[utf8]{{inputenc}}
\\pagestyle{{empty}}
\\definecolor{{primary}}{{HTML}}{{{primary}}}
\\definecolor{{secondary}}{{HTML}}{{{secondary}}}
\\definecolor{{accent}}{{HTML}}{{{accent}}}
\\definecolor{{bgposter}}{{HTML}}{{{bg}}}
\\pagecolor{{bgposter}}
\\setlist[itemize]{{leftmargin=*, itemsep=2mm, topsep=1mm}}
\\begin{{document}}
\\noindent
\\colorbox{{primary}}{{%
\\parbox{{\\dimexpr\\linewidth-0pt\\relax}}{{%
  \\centering\\color{{white}}\\vspace{{6mm}}
  {{\\fontsize{{{title_pt}}}{{{title_pt+12}}}\\selectfont\\bfseries {_latex_escape(title)}}}\\par\\vspace{{4mm}}
  {{\\large {_latex_escape(author)} \\quad|\\quad {_latex_escape(venue)}}}\\par\\vspace{{6mm}}
}}}}
\\vspace{{8mm}}

\\noindent
{stats_row}

\\vspace{{10mm}}
\\noindent
{{\\fontsize{{{body_pt}}}{{{body_pt+8}}}\\selectfont
{body}
}}
\\end{{document}}
"""
    path.write_text(tex, encoding="utf-8")


def _find_latex() -> Optional[str]:
    for name in ("pdflatex", "xelatex"):
        found = shutil.which(name)
        if found:
            return found
    return None


def _compile_pdf(poster_dir: Path, engine: Optional[str] = None) -> Tuple[int, str, str]:
    binary = engine or _find_latex()
    if not binary:
        return 3, "", "pdflatex/xelatex unavailable"
    cmd = [binary, "-interaction=nonstopmode", "-halt-on-error", "main.tex"]
    outs: List[str] = []
    errs: List[str] = []
    rc = 0
    for _ in range(2):
        proc = subprocess.run(
            cmd,
            cwd=str(poster_dir),
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=180,
        )
        outs.append(proc.stdout or "")
        errs.append(proc.stderr or "")
        rc = proc.returncode
        if rc != 0:
            break
    return rc, "\n".join(outs), "\n".join(errs)


def _write_pptx(
    path: Path,
    title: str,
    author: str,
    venue: str,
    size: str,
    orientation: str,
    boxes: Sequence[PosterBox],
    stats: Sequence[str],
    primary: str,
    secondary: str,
    bg: str,
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Mm, Pt

    w, h = SIZE_MM[(size.upper(), orientation.lower())]
    prs = Presentation()
    prs.slide_width = Mm(w)
    prs.slide_height = Mm(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(int(bg[0:2], 16), int(bg[2:4], 16), int(bg[4:6], 16))
    primary_rgb = RGBColor(int(primary[0:2], 16), int(primary[2:4], 16), int(primary[4:6], 16))
    secondary_rgb = RGBColor(int(secondary[0:2], 16), int(secondary[2:4], 16), int(secondary[4:6], 16))

    # title bar
    title_shape = slide.shapes.add_shape(1, Mm(0), Mm(0), Mm(w), Mm(h * 0.14))  # MSO_SHAPE.RECTANGLE = 1
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = primary_rgb
    title_shape.line.fill.background()
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40 if size.upper() == "A0" else 28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"{author} | {venue}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER

    # stats
    stat_y = Mm(h * 0.16)
    stat_h = Mm(h * 0.08)
    for i, s in enumerate((list(stats) + ["—"] * 4)[:4]):
        left = Mm(w * (0.04 + i * 0.24))
        box = slide.shapes.add_textbox(left, stat_y, Mm(w * 0.20), stat_h)
        stf = box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = s
        sp.font.size = Pt(28 if size.upper() == "A0" else 18)
        sp.font.bold = True
        sp.font.color.rgb = primary_rgb
        sp.alignment = PP_ALIGN.CENTER

    # content boxes in grid
    n = max(1, len(boxes))
    cols = 2 if n <= 4 else 3
    rows = (n + cols - 1) // cols
    grid_top = h * 0.28
    grid_h = h * 0.68
    cell_w = w * 0.92 / cols
    cell_h = grid_h / rows
    for idx, box in enumerate(boxes):
        r, c = divmod(idx, cols)
        left = Mm(w * 0.04 + c * cell_w)
        top = Mm(grid_top + r * cell_h)
        shape = slide.shapes.add_textbox(left, top, Mm(cell_w * 0.95), Mm(cell_h * 0.92))
        btf = shape.text_frame
        btf.word_wrap = True
        hp = btf.paragraphs[0]
        hp.text = box.title
        hp.font.size = Pt(24 if size.upper() == "A0" else 16)
        hp.font.bold = True
        hp.font.color.rgb = secondary_rgb
        for bullet in box.bullets:
            bp = btf.add_paragraph()
            bp.text = f"• {bullet}"
            bp.font.size = Pt(16 if size.upper() == "A0" else 12)
            bp.font.color.rgb = RGBColor(0x11, 0x18, 0x27)
    prs.save(str(path))


def _copy_figures(paper_dir: Path, poster_dir: Path) -> List[str]:
    src = paper_dir / "figures"
    dst = poster_dir / "figures"
    dst.mkdir(parents=True, exist_ok=True)
    copied = []
    if not src.is_dir():
        return copied
    for path in src.rglob("*"):
        if path.is_file() and path.suffix.lower() in {".pdf", ".png", ".jpg", ".jpeg", ".svg"}:
            target = dst / path.name
            shutil.copy2(path, target)
            copied.append(target.name)
    return copied


def build(
    workspace: Path,
    *,
    venue: str = "NeurIPS",
    size: str = "A0",
    orientation: str = "landscape",
    columns: int = 4,
    engine: str = "",
) -> dict:
    workspace = Path(workspace).expanduser().resolve()
    paper_dir = workspace / "paper"
    poster_dir = workspace / "poster"
    poster_dir.mkdir(parents=True, exist_ok=True)
    if not paper_dir.is_dir():
        raise SystemExit("paper/ directory is missing")

    size = size.upper()
    orientation = orientation.lower()
    if (size, orientation) not in SIZE_MM:
        size, orientation = "A0", "landscape"
    columns = max(2, min(int(columns or 4), 4))
    if orientation == "portrait":
        columns = min(columns, 3)

    title, author, sections, stats = _extract(paper_dir)
    boxes = _plan_boxes(sections)
    theme = VENUE_THEME.get(venue.strip().lower(), VENUE_THEME["generic"])
    primary, secondary, accent, bg = theme
    figures = _copy_figures(paper_dir, poster_dir)

    _write_plan(poster_dir / "POSTER_CONTENT_PLAN.md", title, venue, size, orientation, columns, boxes, stats)
    _write_speech(poster_dir / "POSTER_SPEECH.md", title, venue, boxes)
    _write_tex(
        poster_dir / "main.tex",
        title,
        author,
        venue,
        size,
        orientation,
        columns,
        boxes,
        stats,
        primary,
        secondary,
        accent,
        bg,
    )
    _write_pptx(
        poster_dir / "poster.pptx",
        title,
        author,
        venue,
        size,
        orientation,
        boxes,
        stats,
        primary,
        secondary,
        bg,
    )
    (poster_dir / "generate_pptx.py").write_text(
        "# Regenerated by host paper-poster builder; re-run build_poster.py to refresh.\n",
        encoding="utf-8",
    )

    rc, stdout, stderr = _compile_pdf(poster_dir, engine or None)
    pdf = poster_dir / "main.pdf"
    state = {
        "phase": 8,
        "venue": venue,
        "poster_size": size,
        "orientation": orientation,
        "columns": columns,
        "status": "completed" if pdf.is_file() and pdf.stat().st_size >= 500 else "partial",
        "timestamp": _now_iso(),
        "figures_selected": figures,
        "executor": "host_step_runner",
        "compile_returncode": rc,
    }
    (poster_dir / "POSTER_STATE.json").write_text(json.dumps(state, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (poster_dir / "POSTER_REVIEW.md").write_text(
        "# Poster Review\n\nHost deterministic builder completed without external reviewer.\n"
        f"- Boxes: {len(boxes)}\n- Compile rc: {rc}\n",
        encoding="utf-8",
    )

    report = {
        "success": pdf.is_file() and pdf.stat().st_size >= 500 and (poster_dir / "poster.pptx").is_file(),
        "returncode": 0 if (pdf.is_file() and pdf.stat().st_size >= 500) else (rc or 4),
        "title": title,
        "artifacts": [
            "poster/main.tex",
            "poster/main.pdf",
            "poster/poster.pptx",
            "poster/POSTER_CONTENT_PLAN.md",
            "poster/POSTER_SPEECH.md",
            "poster/POSTER_STATE.json",
        ],
        "stdout_tail": (stdout or "")[-2000:],
        "stderr_tail": (stderr or "")[-2000:],
        "state": state,
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return report


def main(argv: Optional[Sequence[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Build conference poster from paper/")
    parser.add_argument("--workspace", required=True)
    parser.add_argument("--venue", default="NeurIPS")
    parser.add_argument("--size", default="A0")
    parser.add_argument("--orientation", default="landscape")
    parser.add_argument("--columns", type=int, default=4)
    parser.add_argument("--engine", default="")
    args = parser.parse_args(argv)
    report = build(
        Path(args.workspace),
        venue=args.venue,
        size=args.size,
        orientation=args.orientation,
        columns=args.columns,
        engine=args.engine,
    )
    return 0 if report.get("success") else int(report.get("returncode") or 4)


if __name__ == "__main__":
    raise SystemExit(main())
